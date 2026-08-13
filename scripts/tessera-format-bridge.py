#!/usr/bin/env python3
"""
tessera-format-bridge.py

A single-file format bridge between Tessera's Block AST and external
document formats. Reads a JSON command from stdin, processes the file,
and writes the result to stdout as JSON.

Input format:
    {"action": "import", "format": "docx", "data": "<base64>"}
    {"action": "export", "format": "docx", "blocks": <DocumentAST JSON>}

Output format:
    {"status": "ok", "data": "<base64>", "metadata": {...}}
    {"error": "<message>"}

Supported formats (import / export):
    docx   - python-docx (read) / Pandoc + python-docx (write)
    xlsx   - openpyxl
    pptx   - python-pptx
    eml    - mailbox stdlib
    mbox   - mailbox stdlib
    html   - beautifulsoup4 + Pandoc (import) / self-contained HTML (export)
    md     - Pandoc (import) / self-contained Markdown (export)
    pdf    - pdftotext (import) / weasyprint HTML (export)

Fallback: Pandoc as universal converter for any format Pandoc handles.
"""

from __future__ import annotations

import base64
import json
import os
import shutil
import subprocess
import sys
import tempfile
import uuid
from typing import Any, Optional

# ---------------------------------------------------------------------------
# Format name normalization
# ---------------------------------------------------------------------------

# Map common format names to Pandoc input format names.
_FORMAT_TO_PANDOC_IN: dict[str, str] = {
    "md": "markdown",
    "markdown": "markdown",
    "rst": "rst",
    "rtf": "rtf",
    "odt": "odt",
    "epub": "epub",
    "latex": "latex",
    "org": "org",
    "html": "html",
    "htm": "html",
}


def _pandoc_input_format(fmt: str) -> str:
    """Map a format name to the Pandoc input format name."""
    return _FORMAT_TO_PANDOC_IN.get(fmt.lower(), fmt.lower())


# ---------------------------------------------------------------------------
# Library availability helpers
# ---------------------------------------------------------------------------

def _try_import(name: str):
    try:
        __import__(name)
        return True
    except ImportError:
        return False


_HAVE_PANDOC = bool(shutil.which("pandoc"))
_HAVE_PDFTOTEXT = bool(shutil.which("pdftotext"))
_HAVE_DOCX = _try_import("docx")
_HAVE_OPENPYXL = _try_import("openpyxl")
_HAVE_PPTX = _try_import("pptx")
_HAVE_BS4 = _try_import("bs4") and _try_import("html.parser")
_HAVE_WEASYPRINT = _try_import("weasyprint")


# ---------------------------------------------------------------------------
# Block AST utilities
# ---------------------------------------------------------------------------

def _make_uuid() -> str:
    return str(uuid.uuid4())


def _plain_runs(text: str) -> list[dict]:
    """Wrap plain text as the inline-runs list format."""
    return [{"text": text, "annotations": []}]


def _inline_runs_from_xml_para(para) -> list[dict]:
    """
    Convert an lxml element (paragraph run / text node) to inline runs.
    Supports bold, italic, underline, strike, code, link.
    """
    runs = []
    for child in para.iter():
        tag = child.tag if child.tag is not None else ""
        if hasattr(child, "text") and child.text:
            bold = child.tag == "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}b"
            italic = child.tag == "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}i"
            code = child.tag == "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}code"
            anns = []
            if bold:
                anns.append("bold")
            if italic:
                anns.append("italic")
            if code:
                anns.append("code")
            runs.append({"text": child.text, "annotations": anns})
        if hasattr(child, "tail") and child.tail:
            runs.append({"text": child.tail, "annotations": []})
    if not runs:
        runs.append({"text": "", "annotations": []})
    return runs


def _build_blocks_from_docx_xml(doc_xml_bytes: bytes) -> tuple[dict[str, Any], list[str]]:
    """
    Parse DOCX XML directly (avoids python-docx dependency) and build
    a Block AST. Returns (blocks, rootChildren).
    """
    from lxml import etree

    nsmap = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    }
    body = etree.fromstring(doc_xml_bytes, etree.XMLParser(recover=True))
    blocks: dict[str, Any] = {}
    root_children: list[str] = []

    def q(tag: str) -> str:
        return "{" + nsmap["w"] + "}" + tag

    for para in body.iter(q("p")):
        pStyle = para.get(q("pStyle"))
        run = para.find(q("r"))
        if pStyle and pStyle.startswith("Heading"):
            try:
                level = int(pStyle.replace("Heading", ""))
            except ValueError:
                level = 1
            block_id = _make_uuid()
            text = "".join(r.text or "" for r in para.iter(q("t")))
            blocks[block_id] = {
                "id": block_id,
                "type": "heading",
                "attributes": {"level": level},
                "content": _plain_runs(text),
                "children": [],
                "parentID": None,
            }
            root_children.append(block_id)
        else:
            block_id = _make_uuid()
            text = "".join(r.text or "" for r in para.iter(q("t")))
            blocks[block_id] = {
                "id": block_id,
                "type": "paragraph",
                "attributes": {},
                "content": _plain_runs(text),
                "children": [],
                "parentID": None,
            }
            root_children.append(block_id)

    # Fallback: if no blocks (e.g., empty body), add a placeholder.
    if not blocks:
        block_id = _make_uuid()
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(""),
            "children": [],
            "parentID": None,
        }
        root_children.append(block_id)

    return blocks, root_children


# ---------------------------------------------------------------------------
# Format detectors
# ---------------------------------------------------------------------------

_EXTENSION_TO_FORMAT: dict[str, str] = {
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".eml": "eml",
    ".mbox": "mbox",
    ".html": "html",
    ".htm": "html",
    ".md": "md",
    ".markdown": "md",
    ".pdf": "pdf",
    ".rst": "rst",
    ".rtf": "rtf",
    ".odt": "odt",
    ".epub": "epub",
    ".tex": "latex",
    ".org": "org",
}


def _detect_format(data: bytes, filename_hint: str = "") -> str:
    """Detect format from magic bytes and filename extension."""
    # Try extension first
    for ext, fmt in _EXTENSION_TO_FORMAT.items():
        if filename_hint.lower().endswith(ext):
            return fmt
    # Magic bytes
    if data.startswith(b"PK"):
        # Could be docx/xlsx/pptx (all ZIP-based)
        if b"word/" in data[:65536]:
            return "docx"
        if b"xl/" in data[:65536] or b"xlSharedStrings" in data[:65536]:
            return "xlsx"
        if b"ppt/" in data[:65536]:
            return "pptx"
    if data.startswith(b"%PDF"):
        return "pdf"
    if b"From " in data[:512] and b"MAIL" in data[:512]:
        return "mbox"
    if b"<!doctype html" in data[:256].lower() or b"<html" in data[:256].lower():
        return "html"
    return "unknown"


# ---------------------------------------------------------------------------
# Import: file bytes -> Block AST JSON
# ---------------------------------------------------------------------------

def _import_docx(data: bytes) -> dict[str, Any]:
    """Import DOCX: unzip, parse body XML, build Block AST."""
    import zipfile
    from io import BytesIO

    with zipfile.ZipFile(BytesIO(data)) as z:
        with z.open("word/document.xml") as f:
            xml_bytes = f.read()

    blocks, root_children = _build_blocks_from_docx_xml(xml_bytes)
    return {"blocks": blocks, "rootChildren": root_children}


def _import_xlsx(data: bytes) -> dict[str, Any]:
    """Import XLSX: unzip, extract shared strings and sheet1, build Block AST."""
    import zipfile
    from io import BytesIO

    blocks: dict[str, Any] = {}
    root_children: list[str] = []
    shared_strings: list[str] = []

    with zipfile.ZipFile(BytesIO(data)) as z:
        # Read shared strings
        try:
            with z.open("xl/sharedStrings.xml") as f:
                import xml.etree.ElementTree as ET
                tree = ET.parse(f)
                root = tree.getroot()
                ns = {"ns": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
                for si in root.findall("ns:si", ns):
                    t_nodes = si.findall(".//ns:t", ns)
                    text = "".join(t.text or "" for t in t_nodes)
                    shared_strings.append(text)
        except KeyError:
            pass

        # Read first sheet
        sheet_paths = [p for p in z.namelist() if p.startswith("xl/worksheets/sheet") and p.endswith(".xml")]
        sheet_paths.sort()
        sheet_path = sheet_paths[0] if sheet_paths else None

        if sheet_path:
            import xml.etree.ElementTree as ET
            with z.open(sheet_path) as f:
                tree = ET.parse(f)
                root = tree.getroot()
            ns = {"ns": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
            for row in root.findall(".//ns:row", ns):
                for cell in row.findall("ns:c", ns):
                    cell_ref = cell.get("r", "")
                    t_attr = cell.get("t", "")
                    v = cell.find("ns:v", ns)
                    raw_val = ""
                    if v is not None and v.text:
                        if t_attr == "s":
                            idx = int(v.text)
                            raw_val = shared_strings[idx] if idx < len(shared_strings) else ""
                        else:
                            raw_val = v.text
                    else:
                        raw_val = ""
                    if raw_val:
                        block_id = _make_uuid()
                        blocks[block_id] = {
                            "id": block_id,
                            "type": "paragraph",
                            "attributes": {"cellRef": cell_ref},
                            "content": _plain_runs(raw_val),
                            "children": [],
                            "parentID": None,
                        }
                        root_children.append(block_id)

    if not blocks:
        block_id = _make_uuid()
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(""),
            "children": [],
            "parentID": None,
        }
        root_children.append(block_id)

    return {"blocks": blocks, "rootChildren": root_children}


def _import_pptx(data: bytes) -> dict[str, Any]:
    """Import PPTX: unzip, extract slide text, build Block AST."""
    import zipfile
    from io import BytesIO
    import xml.etree.ElementTree as ET

    blocks: dict[str, Any] = {}
    root_children: list[str] = []

    with zipfile.ZipFile(BytesIO(data)) as z:
        slide_paths = [p for p in z.namelist() if p.startswith("ppt/slides/slide") and p.endswith(".xml")]
        slide_paths.sort()

        for i, slide_path in enumerate(slide_paths, 1):
            with z.open(slide_path) as f:
                tree = ET.parse(f)
                root = tree.getroot()

            # Slide title
            title_id = _make_uuid()
            blocks[title_id] = {
                "id": title_id,
                "type": "heading",
                "attributes": {"level": 1},
                "content": _plain_runs(f"Slide {i}"),
                "children": [],
                "parentID": None,
            }
            root_children.append(title_id)

            # All text in the slide
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            for t in root.iter():
                if hasattr(t, "text") and t.text and t.text.strip():
                    block_id = _make_uuid()
                    blocks[block_id] = {
                        "id": block_id,
                        "type": "paragraph",
                        "attributes": {},
                        "content": _plain_runs(t.text.strip()),
                        "children": [],
                        "parentID": None,
                    }
                    root_children.append(block_id)

    if not blocks:
        block_id = _make_uuid()
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(""),
            "children": [],
            "parentID": None,
        }
        root_children.append(block_id)

    return {"blocks": blocks, "rootChildren": root_children}


def _import_eml_mbox(data: bytes, filename_hint: str = "") -> dict[str, Any]:
    """Import EML or MBOX: parse email messages, build Block AST."""
    import email
    from email.policy import default as email_default

    blocks: dict[str, Any] = {}
    root_children: list[str] = []

    if b"From " in data[:512] and b"MAIL" not in data[:256]:
        # MBOX format: split on From lines
        raw_messages = data.decode("utf-8", errors="replace").split("\nFrom ")
        raw_messages[0] = raw_messages[0].replace("\nFrom ", "", 1) if "\nFrom " in raw_messages[0] else raw_messages[0]
    else:
        # Single EML
        raw_messages = [data.decode("utf-8", errors="replace")]

    for i, raw in enumerate(raw_messages, 1):
        try:
            msg = email.message_from_string(raw if isinstance(raw, str) else raw)
        except Exception:
            block_id = _make_uuid()
            blocks[block_id] = {
                "id": block_id,
                "type": "paragraph",
                "attributes": {},
                "content": _plain_runs(raw[:200] if isinstance(raw, str) else ""),
                "children": [],
                "parentID": None,
            }
            root_children.append(block_id)
            continue

        subject = msg.get("Subject", "(no subject)")
        sender = msg.get("From", "")
        date = msg.get("Date", "")
        body = ""

        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    charset = part.get_content_charset() or "utf-8"
                    body = part.get_payload(decode=True).decode(charset, errors="replace")
                    break
        else:
            charset = msg.get_content_charset() or "utf-8"
            body = msg.get_payload(decode=True).decode(charset, errors="replace")

        # Heading
        heading_id = _make_uuid()
        header_text = f"{subject} | {sender} | {date}"
        blocks[heading_id] = {
            "id": heading_id,
            "type": "heading",
            "attributes": {"level": 2},
            "content": _plain_runs(header_text),
            "children": [],
            "parentID": None,
        }
        root_children.append(heading_id)

        # Body
        body_id = _make_uuid()
        blocks[body_id] = {
            "id": body_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(body.strip()),
            "children": [],
            "parentID": None,
        }
        root_children.append(body_id)

    if not blocks:
        block_id = _make_uuid()
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(""),
            "children": [],
            "parentID": None,
        }
        root_children.append(block_id)

    return {"blocks": blocks, "rootChildren": root_children}


def _import_html(data: bytes) -> dict[str, Any]:
    """Import HTML: parse with beautifulsoup4, convert to Block AST."""
    if not _HAVE_BS4:
        # Fall back to Pandoc
        return _import_pandoc(data, "html")

    from bs4 import BeautifulSoup

    text = BeautifulSoup(data.decode("utf-8", errors="replace"), "html.parser").get_text("\n")
    # Split into paragraphs
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    blocks: dict[str, Any] = {}
    root_children: list[str] = []

    for para in paragraphs:
        block_id = _make_uuid()
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(para),
            "children": [],
            "parentID": None,
        }
        root_children.append(block_id)

    if not blocks:
        block_id = _make_uuid()
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(""),
            "children": [],
            "parentID": None,
        }
        root_children.append(block_id)

    return {"blocks": blocks, "rootChildren": root_children}


def _import_pdf(data: bytes) -> dict[str, Any]:
    """Import PDF: use pdftotext for text extraction."""
    if not _HAVE_PDFTOTEXT:
        return _import_pandoc(data, "pdf")

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        f.write(data)
        tmp_path = f.name

    try:
        result = subprocess.run(
            ["pdftotext", "-layout", tmp_path, "-"],
            capture_output=True,
            text=True,
            timeout=30,
        )
        text = result.stdout
    except (subprocess.TimeoutExpired, FileNotFoundError):
        text = ""
    finally:
        os.unlink(tmp_path)

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    blocks: dict[str, Any] = {}
    root_children: list[str] = []

    for para in paragraphs:
        block_id = _make_uuid()
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(para),
            "children": [],
            "parentID": None,
        }
        root_children.append(block_id)

    if not blocks:
        block_id = _make_uuid()
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": _plain_runs(""),
            "children": [],
            "parentID": None,
        }
        root_children.append(block_id)

    return {"blocks": blocks, "rootChildren": root_children}


def _import_pandoc(data: bytes, format_hint: str = "") -> dict[str, Any]:
    """
    Import via Pandoc: convert to JSON AST, build Block AST preserving
    block types (heading, paragraph, code block, quote, list).
    """
    if not _HAVE_PANDOC:
        raise ImportError("pandoc not found on PATH")

    fmt = _pandoc_input_format(format_hint) if format_hint else "markdown"
    with tempfile.NamedTemporaryFile(suffix=".bin", delete=False) as f:
        f.write(data)
        tmp_in = f.name

    try:
        result = subprocess.run(
            ["pandoc", "-f", fmt, "-t", "json", "--quiet", tmp_in],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode != 0:
            raise RuntimeError(f"pandoc import failed: {result.stderr}")

        blocks: dict[str, Any] = {}
        root_children: list[str] = []

        try:
            pandoc_ast = json.loads(result.stdout)
        except json.JSONDecodeError:
            pandoc_ast = {}

        pandoc_blocks = pandoc_ast.get("blocks", []) if isinstance(pandoc_ast, dict) else []

        for pb in pandoc_blocks:
            block = _build_block_from_pandoc(pb, blocks, root_children)
            if block:
                root_children.append(block)

        if not blocks:
            block_id = _make_uuid()
            blocks[block_id] = {
                "id": block_id,
                "type": "paragraph",
                "attributes": {},
                "content": _plain_runs(""),
                "children": [],
                "parentID": None,
            }
            root_children.append(block_id)

        return {"blocks": blocks, "rootChildren": root_children}
    finally:
        os.unlink(tmp_in)


def _build_block_from_pandoc(
    pb: dict, blocks: dict, root_children: list
) -> Optional[str]:
    """
    Convert a Pandoc JSON block to a Block AST block.
    Returns the block ID if a block was created, None otherwise.
    """
    pt = pb.get("t", "")
    pc = pb.get("c")
    if pc is None:
        pc = []

    block_id = _make_uuid()

    def extract_text(content) -> str:
        """Walk Pandoc inline content and extract plain text."""
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            return "".join(extract_text(item) for item in content)
        if isinstance(content, dict):
            t = content.get("t", "")
            c = content.get("c")
            if t == "Str" and isinstance(c, str):
                return c
            if t in ("Space", "SoftBreak"):
                return " "
            if t == "LineBreak":
                return "\n"
            if c is not None:
                return extract_text(c)
        return ""

    def extract_runs(content) -> list[dict]:
        """Walk Pandoc inline content and build inline runs."""
        runs: list[dict] = []
        anns: list[str] = []

        def walk(node):
            if isinstance(node, list):
                for item in node:
                    walk(item)
            elif isinstance(node, dict):
                t = node.get("t", "")
                c = node.get("c")
                if t == "Str" and isinstance(c, str):
                    runs.append({"text": c, "annotations": list(anns)})
                elif t == "Emph":
                    anns.append("italic")
                    if c is not None:
                        walk(c)
                    anns.pop()
                elif t == "Strong":
                    anns.append("bold")
                    if c is not None:
                        walk(c)
                    anns.pop()
                elif t == "Strikeout":
                    anns.append("strikethrough")
                    if c is not None:
                        walk(c)
                    anns.pop()
                elif t == "Code":
                    anns.append("code")
                    text = c.get("text", "") if isinstance(c, dict) else str(c)
                    runs.append({"text": text, "annotations": list(anns)})
                    anns.pop()
                elif t == "Link":
                    anns.append("link")
                    if c is not None:
                        walk(c)
                    if anns and anns[-1] == "link":
                        anns.pop()
                elif t == "Space":
                    runs.append({"text": " ", "annotations": list(anns)})
                elif t == "SoftBreak":
                    runs.append({"text": " ", "annotations": list(anns)})
                elif c is not None:
                    walk(c)

        walk(content)
        return runs if runs else [{"text": "", "annotations": []}]

    if pt == "Header":
        # c = [level, identifier+classes+attributes, inline_content]
        level = 1
        content_list = pc
        if isinstance(pc, list) and len(pc) >= 1:
            level = int(pc[0]) if isinstance(pc[0], (int, float)) else 1
        if isinstance(pc, list) and len(pc) >= 3:
            content_list = pc[2]
        blocks[block_id] = {
            "id": block_id,
            "type": "heading",
            "attributes": {"level": level},
            "content": extract_runs(content_list),
            "children": [],
            "parentID": None,
        }
        return block_id

    elif pt == "Para":
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": extract_runs(pc),
            "children": [],
            "parentID": None,
        }
        return block_id

    elif pt == "Plain":
        blocks[block_id] = {
            "id": block_id,
            "type": "paragraph",
            "attributes": {},
            "content": extract_runs(pc),
            "children": [],
            "parentID": None,
        }
        return block_id

    elif pt == "CodeBlock":
        # c = [identifier, classes, attributes, text]
        text = ""
        if isinstance(pc, list) and len(pc) >= 4:
            text = str(pc[3])
        elif isinstance(pc, str):
            text = pc
        language = ""
        if isinstance(pc, list) and len(pc) >= 2:
            classes = pc[1]
            if isinstance(classes, list) and classes:
                language = classes[0] if classes else ""
        blocks[block_id] = {
            "id": block_id,
            "type": "codeBlock",
            "attributes": {"language": language},
            "content": _plain_runs(text),
            "children": [],
            "parentID": None,
        }
        return block_id

    elif pt == "BlockQuote":
        inner_blocks: list[str] = []
        inner_content = pc if isinstance(pc, list) else [pc]
        for inner_pb in inner_content:
            if isinstance(inner_pb, dict):
                inner_id = _build_block_from_pandoc(inner_pb, blocks, inner_blocks)
                if inner_id:
                    inner_blocks.append(inner_id)
        blocks[block_id] = {
            "id": block_id,
            "type": "quote",
            "attributes": {},
            "content": _plain_runs(extract_text(pc)),
            "children": inner_blocks,
            "parentID": None,
        }
        return block_id

    elif pt in ("BulletList", "OrderedList"):
        list_items = pc[1] if isinstance(pc, list) and len(pc) > 1 else (pc if isinstance(pc, list) else [])
        list_block_id = _make_uuid()
        item_block_ids: list[str] = []
        for item_content in list_items:
            item_block_id = _make_uuid()
            item_runs = extract_runs(item_content) if isinstance(item_content, list) else []
            blocks[item_block_id] = {
                "id": item_block_id,
                "type": "paragraph",
                "attributes": {},
                "content": item_runs,
                "children": [],
                "parentID": list_block_id,
            }
            item_block_ids.append(item_block_id)
        blocks[list_block_id] = {
            "id": list_block_id,
            "type": "list",
            "attributes": {
                "style": "unordered" if pt == "BulletList" else "ordered"
            },
            "content": [],
            "children": item_block_ids,
            "parentID": None,
        }
        # Don't add the list_block_id to root_children here;
        # the original pb was the list itself
        return list_block_id

    elif pt == "HorizontalRule":
        blocks[block_id] = {
            "id": block_id,
            "type": "divider",
            "attributes": {},
            "content": [],
            "children": [],
            "parentID": None,
        }
        return block_id

    else:
        # Unknown block type: emit as paragraph with plain text
        text = extract_text(pc)
        if text.strip():
            blocks[block_id] = {
                "id": block_id,
                "type": "paragraph",
                "attributes": {},
                "content": _plain_runs(text),
                "children": [],
                "parentID": None,
            }
            return block_id
        return None


def _extract_pandoc_text(ast: Any) -> list[str]:
    """
    Walk a Pandoc JSON AST and extract plain text lines.

    The Pandoc JSON AST root is a dict: {"pandoc-api-version": [...],
    "meta": {...}, "blocks": [...]}. Each block is
    {"t": "BlockType", "c": content}. Inline elements within
    block content are {"t": "InlineType", "c": content}.
    """
    lines: list[str] = []

    def walk_inline(node):
        """Extract text from an inline element."""
        if isinstance(node, list):
            for item in node:
                walk_inline(item)
        elif isinstance(node, dict):
            t = node.get("t", "")
            c = node.get("c")
            if t == "Str" and isinstance(c, str):
                lines.append(c)
            elif t == "Space":
                lines.append(" ")
            elif t == "SoftBreak":
                lines.append(" ")
            elif t == "LineBreak":
                lines.append("\n")
            elif c is not None:
                walk_inline(c)

    def walk_block(node):
        """Extract text from a block element."""
        if isinstance(node, list):
            for item in node:
                walk_block(item)
        elif isinstance(node, dict):
            t = node.get("t", "")
            c = node.get("c")
            # These block types carry inline content
            if t in ("Para", "Header", "CodeBlock", "BlockQuote", "Plain", "Div"):
                if c is not None:
                    walk_inline(c)
                lines.append("\n")
            elif c is not None:
                walk_block(c)

    # Pandoc JSON AST: {"blocks": [...], ...}
    if isinstance(ast, dict):
        blocks = ast.get("blocks", [])
        walk_block(blocks)
    elif isinstance(ast, list):
        walk_block(ast)

    # Join and split into paragraphs
    text = "".join(lines)
    return [p.strip() for p in text.split("\n") if p.strip()]


def _do_import(format_hint: str, data: bytes) -> dict[str, Any]:
    """Route to the correct import handler."""
    fmt = format_hint.lower()

    # If format is unknown, detect
    if fmt == "unknown":
        fmt = _detect_format(data)

    handlers: dict[str, callable] = {
        "docx": _import_docx,
        "xlsx": _import_xlsx,
        "pptx": _import_pptx,
        "eml": lambda d: _import_eml_mbox(d, "eml"),
        "mbox": lambda d: _import_eml_mbox(d, "mbox"),
        "html": _import_html,
        "pdf": _import_pdf,
    }

    if fmt in handlers:
        return handlers[fmt](data)

    # Try Pandoc for anything Pandoc handles
    if _HAVE_PANDOC:
        return _import_pandoc(data, fmt if fmt != "unknown" else "")

    raise ImportError(f"no converter for .{fmt}")


# ---------------------------------------------------------------------------
# Export: Block AST JSON -> file bytes
# ---------------------------------------------------------------------------

def _export_markdown(blocks: dict[str, Any]) -> bytes:
    """Export Block AST to Markdown."""
    root_children = blocks.get("rootChildren", [])
    block_map = blocks.get("blocks", {})

    lines: list[str] = []

    for block_id in root_children:
        block = block_map.get(block_id, {})
        btype = block.get("type", "")
        content_runs = block.get("content", [])
        text = "".join(r.get("text", "") for r in content_runs)

        if btype == "heading":
            level = block.get("attributes", {}).get("level", 1)
            lines.append(f"{'#' * level} {text}")
        elif btype == "paragraph":
            # Reconstruct inline formatting
            out = ""
            for run in content_runs:
                t = run.get("text", "")
                anns = run.get("annotations", [])
                if "bold" in anns:
                    t = f"**{t}**"
                if "italic" in anns:
                    t = f"*{t}*"
                if "code" in anns:
                    t = f"`{t}`"
                out += t
            lines.append(out)
        elif btype == "codeBlock":
            lang = block.get("attributes", {}).get("language", "")
            lines.append(f"```{lang}\n{text}\n```")
        elif btype == "quote":
            for ln in text.splitlines():
                lines.append(f"> {ln}")
        elif btype == "list":
            items = block.get("children", [])
            for i, item_id in enumerate(items, 1):
                item = block_map.get(item_id, {})
                item_text = "".join(r.get("text", "") for r in item.get("content", []))
                style = block.get("attributes", {}).get("style", "unordered")
                prefix = f"{i}." if style == "ordered" else "-"
                lines.append(f"{prefix} {item_text}")
        lines.append("")

    return "\n".join(lines).encode("utf-8")


def _export_html(blocks: dict[str, Any]) -> bytes:
    """Export Block AST to self-contained HTML."""
    root_children = blocks.get("rootChildren", [])
    block_map = blocks.get("blocks", {})

    body_parts: list[str] = []

    for block_id in root_children:
        block = block_map.get(block_id, {})
        btype = block.get("type", "")
        content_runs = block.get("content", [])
        text = "".join(r.get("text", "") for r in content_runs)
        escape = lambda s: s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        def render_runs(runs):
            out = ""
            for run in runs:
                t = escape(run.get("text", ""))
                anns = run.get("annotations", [])
                if "bold" in anns:
                    t = f"<strong>{t}</strong>"
                if "italic" in anns:
                    t = f"<em>{t}</em>"
                if "code" in anns:
                    t = f"<code>{t}</code>"
                out += t
            return out

        if btype == "heading":
            level = block.get("attributes", {}).get("level", 1)
            body_parts.append(f"<h{level}>{render_runs(content_runs)}</h{level}>")
        elif btype == "paragraph":
            body_parts.append(f"<p>{render_runs(content_runs)}</p>")
        elif btype == "codeBlock":
            lang = escape(block.get("attributes", {}).get("language", ""))
            body_parts.append(f'<pre><code class="language-{lang}">{escape(text)}</code></pre>')
        elif btype == "quote":
            body_parts.append(f"<blockquote><p>{render_runs(content_runs)}</p></blockquote>")
        elif btype == "list":
            tag = "ol" if block.get("attributes", {}).get("style") == "ordered" else "ul"
            items = block.get("children", [])
            inner = ""
            for item_id in items:
                item = block_map.get(item_id, {})
                item_text = render_runs(item.get("content", []))
                inner += f"<li>{item_text}</li>"
            body_parts.append(f"<{tag}>{inner}</{tag}>")
        else:
            body_parts.append(f"<p>{render_runs(content_runs)}</p>")

    html = (
        "<!doctype html>\n"
        "<html lang=\"en\">\n"
        "<head>\n<meta charset=\"utf-8\">\n"
        "<title>Tessera Export</title>\n"
        "</head>\n"
        "<body>\n"
        + "\n".join(body_parts)
        + "\n</body>\n</html>\n"
    )
    return html.encode("utf-8")


def _export_pdf(blocks: dict[str, Any]) -> bytes:
    """Export Block AST to PDF via weasyprint (HTML -> PDF)."""
    if not _HAVE_WEASYPRINT:
        raise ImportError("weasyprint not installed")
    from weasyprint import HTML
    html_bytes = _export_html(blocks)
    html_str = html_bytes.decode("utf-8")
    doc = HTML(string=html_str)
    return doc.write_pdf()


def _export_docx(blocks: dict[str, Any]) -> bytes:
    """
    Export Block AST to DOCX.
    Prefers python-docx when available; falls back to Pandoc.
    """
    if _HAVE_DOCX:
        return _export_docx_python(blocks)

    if _HAVE_PANDOC:
        # Write HTML to temp, convert via Pandoc
        html_bytes = _export_html(blocks)
        with tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="wb") as f:
            f.write(html_bytes)
            tmp_in = f.name
        try:
            result = subprocess.run(
                ["pandoc", "-f", "html", "-t", "docx", "--quiet", "-o", "-", tmp_in],
                capture_output=True,
                timeout=30,
            )
            if result.returncode != 0:
                raise RuntimeError(f"pandoc docx export failed: {result.stderr.decode()}")
            return result.stdout
        finally:
            os.unlink(tmp_in)

    raise ImportError("no DOCX writer: install python-docx or pandoc")


def _export_docx_python(blocks: dict[str, Any]) -> bytes:
    """Export Block AST to DOCX using python-docx."""
    import zipfile
    from io import BytesIO
    from docx import Document
    from docx.shared import Pt

    root_children = blocks.get("rootChildren", [])
    block_map = blocks.get("blocks", {})

    doc = Document()

    def add_runs(paragraph, runs):
        for run in runs:
            text = run.get("text", "")
            anns = run.get("annotations", [])
            r = paragraph.add_run(text)
            if "bold" in anns:
                r.bold = True
            if "italic" in anns:
                r.italic = True
            if "code" in anns:
                r.font.name = "Courier New"
                r.font.size = Pt(10)

    for block_id in root_children:
        block = block_map.get(block_id, {})
        btype = block.get("type", "")
        content_runs = block.get("content", [])

        if btype == "heading":
            level = block.get("attributes", {}).get("level", 1)
            p = doc.add_heading("", level=level)
            add_runs(p, content_runs)
        elif btype == "paragraph":
            p = doc.add_paragraph()
            add_runs(p, content_runs)
        elif btype == "codeBlock":
            code_text = "".join(r.get("text", "") for r in content_runs)
            p = doc.add_paragraph()
            r = p.add_run(code_text)
            r.font.name = "Courier New"
            r.font.size = Pt(10)
        elif btype == "quote":
            p = doc.add_paragraph(style="Quote")
            add_runs(p, content_runs)
        elif btype == "list":
            style = block.get("attributes", {}).get("style", "unordered")
            list_tag = "add_numbering" if style == "ordered" else "add_list"
            items = block.get("children", [])
            for item_id in items:
                item = block_map.get(item_id, {})
                item_runs = item.get("content", [])
                p = doc.add_paragraph(style="List Bullet" if style != "ordered" else "List Number")
                add_runs(p, item_runs)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _export_xlsx(blocks: dict[str, Any]) -> bytes:
    """Export Block AST to XLSX using openpyxl."""
    if not _HAVE_OPENPYXL:
        raise ImportError("openpyxl not installed")
    from openpyxl import Workbook

    root_children = blocks.get("rootChildren", [])
    block_map = blocks.get("blocks", {})

    wb = Workbook()
    ws = wb.active
    ws.title = "Tessera Export"

    row = 1
    for block_id in root_children:
        block = block_map.get(block_id, {})
        btype = block.get("type", "")
        content_runs = block.get("content", [])
        text = "".join(r.get("text", "") for r in content_runs)

        if btype == "heading":
            ws.cell(row=row, column=1, value=f"[H{block.get('attributes', {}).get('level', 1)}] {text}")
            row += 1
        elif btype == "paragraph":
            ws.cell(row=row, column=1, value=text)
            row += 1
        elif btype == "list":
            items = block.get("children", [])
            for item_id in items:
                item = block_map.get(item_id, {})
                item_text = "".join(r.get("text", "") for r in item.get("content", []))
                ws.cell(row=row, column=1, value=f"  - {item_text}")
                row += 1
        else:
            ws.cell(row=row, column=1, value=text)
            row += 1

    from io import BytesIO
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _export_pptx(blocks: dict[str, Any]) -> bytes:
    """Export Block AST to PPTX using python-pptx."""
    if not _HAVE_PPTX:
        raise ImportError("python-pptx not installed")
    from pptx import Presentation
    from pptx.util import Pt

    root_children = blocks.get("rootChildren", [])
    block_map = blocks.get("blocks", {})

    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000

    # Collect slide content
    slides_content: list[list[tuple[str, str]]] = []
    current_slide: list[tuple[str, str]] = []
    SLIDE_THRESHOLD = 10  # New slide every 10 blocks

    for block_id in root_children:
        block = block_map.get(block_id, {})
        btype = block.get("type", "")
        content_runs = block.get("content", [])
        text = "".join(r.get("text", "") for r in content_runs)

        if btype == "heading" and block.get("attributes", {}).get("level", 1) == 1:
            if current_slide:
                slides_content.append(current_slide)
            current_slide = [("title", text)]
        elif len(current_slide) >= SLIDE_THRESHOLD and btype == "heading":
            slides_content.append(current_slide)
            current_slide = [("body", text)]
        else:
            current_slide.append(("body", text))

    if current_slide:
        slides_content.append(current_slide)

    if not slides_content:
        slides_content = [[("body", "Empty document")]]

    for slide_content in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_text = next((text for kind, text in slide_content if kind == "title"), "")
        body_texts = [text for kind, text in slide_content if kind == "body"]
        if title_text:
            slide.shapes.title.text = title_text
        if len(slide_content) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            for i, bt in enumerate(body_texts):
                p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
                p.text = bt

    from io import BytesIO
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _export_eml(blocks: dict[str, Any]) -> bytes:
    """Export Block AST to EML (plain-text email)."""
    import email.message

    root_children = blocks.get("rootChildren", [])
    block_map = blocks.get("blocks", {})

    body_lines: list[str] = []
    for block_id in root_children:
        block = block_map.get(block_id, {})
        text = "".join(r.get("text", "") for r in block.get("content", []))
        body_lines.append(text)

    msg = email.message.EmailMessage()
    msg["Subject"] = "Tessera Export"
    msg["From"] = "Tessera <tessera@local>"
    msg["Date"] = email.utils.formatdate()
    msg.set_content("\n\n".join(body_lines))
    return msg.as_bytes()


def _export_pandoc(blocks: dict[str, Any], target_format: str) -> bytes:
    """Export Block AST to any Pandoc-supported format."""
    if not _HAVE_PANDOC:
        raise ImportError("pandoc not installed")

    # Write HTML, convert via Pandoc
    html_bytes = _export_html(blocks)
    with tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="wb") as f:
        f.write(html_bytes)
        tmp_in = f.name

    try:
        result = subprocess.run(
            ["pandoc", "-f", "html", "-t", target_format, "--quiet", "-o", "-", tmp_in],
            capture_output=True,
            timeout=30,
        )
        if result.returncode != 0:
            raise RuntimeError(f"pandoc export failed: {result.stderr.decode()}")
        return result.stdout
    finally:
        os.unlink(tmp_in)


def _do_export(format_hint: str, blocks: dict[str, Any]) -> bytes:
    """Route to the correct export handler."""
    fmt = format_hint.lower()

    handlers: dict[str, callable] = {
        "md": _export_markdown,
        "html": _export_html,
        "pdf": _export_pdf,
        "docx": _export_docx,
        "xlsx": _export_xlsx,
        "pptx": _export_pptx,
        "eml": _export_eml,
    }

    if fmt in handlers:
        return handlers[fmt](blocks)

    # Try Pandoc for anything else Pandoc handles
    return _export_pandoc(blocks, fmt)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def _metadata_for(blocks: dict[str, Any]) -> dict[str, Any]:
    """Build metadata dict from Block AST."""
    block_list = blocks.get("blocks", {})
    root_children = blocks.get("rootChildren", [])
    word_count = 0
    block_count = len(block_list)

    for block in block_list.values():
        text = "".join(r.get("text", "") for r in block.get("content", []))
        word_count += len(text.split())

    return {
        "blockCount": block_count,
        "rootChildCount": len(root_children),
        "wordCount": word_count,
    }


def main() -> None:
    """Read JSON command from stdin, process, write JSON to stdout."""
    try:
        raw = sys.stdin.read()
        if not raw.strip():
            print(json.dumps({"error": "no input"}), file=sys.stderr)
            sys.exit(1)

        cmd = json.loads(raw)
        action = cmd.get("action", "")
        fmt = cmd.get("format", "")
        metadata: Optional[dict[str, Any]] = None

        if action == "import":
            data_b64 = cmd.get("data", "")
            if isinstance(data_b64, list):
                # Already decoded
                data = bytes(data_b64)
            else:
                data = base64.b64decode(data_b64)

            ast = _do_import(fmt, data)
            metadata = _metadata_for(ast)
            # Re-encode AST as JSON
            out_data = json.dumps(ast).encode("utf-8")
            out_b64 = base64.b64encode(out_data).decode("ascii")
            result = {"status": "ok", "data": out_b64, "metadata": metadata}

        elif action == "export":
            blocks = cmd.get("blocks", {})
            if isinstance(blocks, str):
                blocks = json.loads(blocks)

            out_bytes = _do_export(fmt, blocks)
            out_b64 = base64.b64encode(out_bytes).decode("ascii")
            metadata = _metadata_for(blocks)
            result = {"status": "ok", "data": out_b64, "metadata": metadata}

        else:
            result = {"error": f"unknown action '{action}'"}

    except ImportError as e:
        result = {"error": str(e)}
    except json.JSONDecodeError as e:
        result = {"error": f"JSON parse error: {e}"}
    except Exception as e:
        result = {"error": f"{type(e).__name__}: {e}"}

    # Write result to stdout
    print(json.dumps(result))


if __name__ == "__main__":
    main()
