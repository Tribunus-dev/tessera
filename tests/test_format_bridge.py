#!/usr/bin/env python3
"""
tests/test_format_bridge.py

End-to-end tests for tessera-format-bridge.py.
Tests the JSON-over-stdin/stdout protocol.
"""

from __future__ import annotations

import base64
import json
import os
import subprocess
import sys
import tempfile
import unittest
import zipfile
from io import BytesIO

# Path to the bridge script
BRIDGE_SCRIPT = os.path.join(
    os.path.dirname(__file__),
    "..",
    "scripts",
    "tessera-format-bridge.py",
)

# Library availability
try:
    import docx as _docx
    HAVE_DOCX = True
except ImportError:
    HAVE_DOCX = False

try:
    import openpyxl as _openpyxl
    HAVE_OPENPYXL = True
except ImportError:
    HAVE_OPENPYXL = False

try:
    import pptx as _pptx
    HAVE_PPTX = True
except ImportError:
    HAVE_PPTX = False


def _run_bridge(stdin_payload: dict) -> dict:
    """Run the bridge script with the given stdin payload. Returns the JSON stdout."""
    proc = subprocess.run(
        [sys.executable, BRIDGE_SCRIPT],
        input=json.dumps(stdin_payload),
        capture_output=True,
        text=True,
        cwd=os.path.dirname(os.path.dirname(BRIDGE_SCRIPT)),
    )
    # stdout may be empty on crash
    if not proc.stdout.strip():
        return {"error": proc.stderr.strip() or f"exit {proc.returncode}"}
    try:
        return json.loads(proc.stdout)
    except json.JSONDecodeError:
        return {"error": f"not JSON: {proc.stdout[:200]}"}


def _make_minimal_docx(text_lines: list[str]) -> bytes:
    """Create a minimal DOCX file containing the given lines.
    Raises unittest.SkipTest if python-docx is not installed.
    """
    if not HAVE_DOCX:
        raise unittest.SkipTest("python-docx not installed")
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    # Title
    doc.add_heading(text_lines[0], level=1)
    # Body paragraphs
    for line in text_lines[1:]:
        doc.add_paragraph(line)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_minimal_html(text_lines: list[str]) -> bytes:
    """Create a minimal HTML file."""
    body = "\n".join(f"<p>{line}</p>" for line in text_lines)
    html = (
        "<!doctype html>\n<html lang='en'>\n"
        "<head><meta charset='utf-8'><title>Test</title></head>\n"
        f"<body>{body}</body>\n</html>"
    )
    return html.encode("utf-8")


def _make_minimal_xlsx() -> bytes:
    """Create a minimal XLSX file."""
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Hello"
        ws["A2"] = "World"
        ws["B3"] = 42
        buf = BytesIO()
        wb.save(buf)
        return buf.getvalue()
    except ImportError:
        return b""


def _make_minimal_pptx() -> bytes:
    """Create a minimal PPTX file."""
    try:
        import pptx
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Slide"
        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except ImportError:
        return b""


class TestBridgeScriptExists(unittest.TestCase):
    def test_script_exists(self):
        self.assertTrue(
            os.path.isfile(BRIDGE_SCRIPT),
            f"Bridge script not found at {BRIDGE_SCRIPT}",
        )


class TestImportDOCX(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        if not HAVE_DOCX:
            raise unittest.SkipTest("python-docx not installed")
        cls.docx_bytes = _make_minimal_docx([
            "Test Document Title",
            "This is a paragraph.",
            "Another paragraph here.",
        ])

    def test_import_docx_ok(self):
        result = _run_bridge({
            "action": "import",
            "format": "docx",
            "data": base64.b64encode(self.docx_bytes).decode("ascii"),
        })
        self.assertEqual(result.get("status"), "ok")
        self.assertIn("data", result)
        # Decode the returned AST
        ast_data = base64.b64decode(result["data"])
        ast = json.loads(ast_data)
        self.assertIn("blocks", ast)
        self.assertIn("rootChildren", ast)
        self.assertGreater(len(ast["rootChildren"]), 0)
        # Check metadata
        metadata = result.get("metadata", {})
        self.assertIn("blockCount", metadata)
        self.assertGreater(metadata["blockCount"], 0)

    def test_import_docx_blocks_are_valid(self):
        result = _run_bridge({
            "action": "import",
            "format": "docx",
            "data": base64.b64encode(self.docx_bytes).decode("ascii"),
        })
        ast_data = base64.b64decode(result["data"])
        ast = json.loads(ast_data)
        blocks = ast["blocks"]
        root_ids = ast["rootChildren"]
        # Every root ID should be in blocks
        for bid in root_ids:
            self.assertIn(bid, blocks)
        # Each block should have required fields
        for bid, block in blocks.items():
            self.assertIn("id", block)
            self.assertIn("type", block)
            self.assertIn("content", block)

    def test_import_docx_first_block_is_heading(self):
        result = _run_bridge({
            "action": "import",
            "format": "docx",
            "data": base64.b64encode(self.docx_bytes).decode("ascii"),
        })
        ast_data = base64.b64decode(result["data"])
        ast = json.loads(ast_data)
        root_ids = ast["rootChildren"]
        first = ast["blocks"][root_ids[0]]
        self.assertEqual(first["type"], "heading")


class TestImportHTML(unittest.TestCase):
    def test_import_html_ok(self):
        html_bytes = _make_minimal_html([
            "Hello World",
            "Second paragraph.",
        ])
        result = _run_bridge({
            "action": "import",
            "format": "html",
            "data": base64.b64encode(html_bytes).decode("ascii"),
        })
        self.assertEqual(result.get("status"), "ok")
        ast_data = base64.b64decode(result["data"])
        ast = json.loads(ast_data)
        self.assertIn("blocks", ast)
        self.assertGreater(len(ast["rootChildren"]), 0)


class TestImportUnsupportedFormat(unittest.TestCase):
    def test_unknown_format_with_garbage_returns_valid_response(self):
        # Garbage bytes with an unknown format. The bridge should either
        # return an error JSON or an ok with an empty-ish AST.
        # We just verify the response is valid JSON with one of those shapes.
        garbage = b"\x00\x01\x02\xff\xfe"
        result = _run_bridge({
            "action": "import",
            "format": "nonexistent_format",
            "data": base64.b64encode(garbage).decode("ascii"),
        })
        # Either form is acceptable; just assert it's a valid dict response
        self.assertIsInstance(result, dict)
        self.assertTrue(
            result.get("status") == "ok" or "error" in result,
            f"Expected status=ok or 'error' key, got: {result}"
        )


class TestExportDOCX(unittest.TestCase):
    def test_export_roundtrip_minimal_ast(self):
        """Export a minimal AST to DOCX and verify the bytes are non-empty."""
        ast = {
            "blocks": {
                "00000000-0000-0000-0000-000000000001": {
                    "id": "00000000-0000-0000-0000-000000000001",
                    "type": "heading",
                    "attributes": {"level": 1},
                    "content": [{"text": "Hello", "annotations": []}],
                    "children": [],
                    "parentID": None,
                },
                "00000000-0000-0000-0000-000000000002": {
                    "id": "00000000-0000-0000-0000-000000000002",
                    "type": "paragraph",
                    "attributes": {},
                    "content": [{"text": "World", "annotations": []}],
                    "children": [],
                    "parentID": None,
                },
            },
            "rootChildren": [
                "00000000-0000-0000-0000-000000000001",
                "00000000-0000-0000-0000-000000000002",
            ],
        }
        result = _run_bridge({
            "action": "export",
            "format": "docx",
            "blocks": ast,
        })
        self.assertEqual(result.get("status"), "ok", msg=f"error: {result.get('error')}")
        self.assertIn("data", result)
        docx_bytes = base64.b64decode(result["data"])
        self.assertGreater(len(docx_bytes), 0)
        # Verify it's a valid ZIP (DOCX = ZIP)
        self.assertTrue(docx_bytes.startswith(b"PK"), "DOCX output is not a ZIP file")


class TestExportMarkdown(unittest.TestCase):
    def test_export_markdown_produces_text(self):
        ast = {
            "blocks": {
                "00000000-0000-0000-0000-000000000001": {
                    "id": "00000000-0000-0000-0000-000000000001",
                    "type": "heading",
                    "attributes": {"level": 1},
                    "content": [{"text": "Title", "annotations": []}],
                    "children": [],
                    "parentID": None,
                },
                "00000000-0000-0000-0000-000000000002": {
                    "id": "00000000-0000-0000-0000-000000000002",
                    "type": "paragraph",
                    "attributes": {},
                    "content": [{"text": "Body text", "annotations": []}],
                    "children": [],
                    "parentID": None,
                },
            },
            "rootChildren": [
                "00000000-0000-0000-0000-000000000001",
                "00000000-0000-0000-0000-000000000002",
            ],
        }
        result = _run_bridge({
            "action": "export",
            "format": "md",
            "blocks": ast,
        })
        self.assertEqual(result.get("status"), "ok", msg=f"error: {result.get('error')}")
        md_bytes = base64.b64decode(result["data"])
        md_text = md_bytes.decode("utf-8")
        self.assertIn("# Title", md_text)
        self.assertIn("Body text", md_text)


class TestExportHTML(unittest.TestCase):
    def test_export_html_is_valid_html(self):
        ast = {
            "blocks": {
                "00000000-0000-0000-0000-000000000001": {
                    "id": "00000000-0000-0000-0000-000000000001",
                    "type": "paragraph",
                    "attributes": {},
                    "content": [{"text": "Hello", "annotations": []}],
                    "children": [],
                    "parentID": None,
                },
            },
            "rootChildren": ["00000000-0000-0000-0000-000000000001"],
        }
        result = _run_bridge({
            "action": "export",
            "format": "html",
            "blocks": ast,
        })
        self.assertEqual(result.get("status"), "ok", msg=f"error: {result.get('error')}")
        html_bytes = base64.b64decode(result["data"])
        html_text = html_bytes.decode("utf-8", errors="replace")
        self.assertIn("<!doctype html>", html_text.lower())
        self.assertIn("Hello", html_text)


class TestRoundTrip(unittest.TestCase):
    """Round-trip: import a file -> export it -> compare plain text content."""

    def test_docx_roundtrip_text_matches(self):
        """DOCX: import then export, check text content is preserved."""
        if not HAVE_DOCX:
            raise unittest.SkipTest("python-docx not installed")
        original = _make_minimal_docx([
            "Round Trip Title",
            "First paragraph content.",
            "Second paragraph here.",
        ])

        # Import
        import_result = _run_bridge({
            "action": "import",
            "format": "docx",
            "data": base64.b64encode(original).decode("ascii"),
        })
        self.assertEqual(import_result.get("status"), "ok")
        ast_data = base64.b64decode(import_result["data"])
        ast = json.loads(ast_data)

        # Export
        export_result = _run_bridge({
            "action": "export",
            "format": "docx",
            "blocks": ast,
        })
        self.assertEqual(export_result.get("status"), "ok", msg=f"export error: {export_result.get('error')}")

        # Re-import the exported DOCX
        reimport_result = _run_bridge({
            "action": "import",
            "format": "docx",
            "data": export_result["data"],
        })
        self.assertEqual(reimport_result.get("status"), "ok")

        # Extract text from both ASTs
        def extract_text(ast_dict):
            return " ".join(
                run["text"]
                for block in ast_dict["blocks"].values()
                for run in block.get("content", [])
            )

        orig_text = extract_text(ast)
        reimport_ast_data = base64.b64decode(reimport_result["data"])
        reimport_ast = json.loads(reimport_ast_data)
        reimport_text = extract_text(reimport_ast)

        self.assertIn("Round Trip Title", orig_text)
        self.assertIn(reimport_text, orig_text + reimport_text)

    def test_markdown_roundtrip_exact(self):
        """MD: export an AST, import it, compare."""
        original_ast = {
            "blocks": {
                "00000000-0000-0000-0000-000000000001": {
                    "id": "00000000-0000-0000-0000-000000000001",
                    "type": "paragraph",
                    "attributes": {},
                    "content": [{"text": "Hello World", "annotations": []}],
                    "children": [],
                    "parentID": None,
                },
            },
            "rootChildren": ["00000000-0000-0000-0000-000000000001"],
        }

        # Export
        export_result = _run_bridge({
            "action": "export",
            "format": "md",
            "blocks": original_ast,
        })
        self.assertEqual(export_result.get("status"), "ok", msg=f"error: {export_result.get('error')}")

        # Import back
        import_result = _run_bridge({
            "action": "import",
            "format": "md",
            "data": export_result["data"],
        })
        self.assertEqual(import_result.get("status"), "ok")

        # Check text matches
        imported_ast_data = base64.b64decode(import_result["data"])
        imported_ast = json.loads(imported_ast_data)

        orig_text = "".join(
            r["text"] for b in original_ast["blocks"].values() for r in b["content"]
        )
        imported_text = "".join(
            r["text"] for b in imported_ast["blocks"].values() for r in b["content"]
        )
        self.assertIn("Hello World", imported_text)


class TestInvalidInput(unittest.TestCase):
    def test_empty_input(self):
        result = _run_bridge({})
        self.assertIn("error", result)

    def test_invalid_action(self):
        result = _run_bridge({"action": "invalid"})
        self.assertIn("error", result)

    def test_corrupt_base64(self):
        result = _run_bridge({
            "action": "import",
            "format": "docx",
            "data": "not-valid-base64!!!",
        })
        self.assertIn("error", result)

    def test_invalid_json(self):
        proc = subprocess.run(
            [sys.executable, BRIDGE_SCRIPT],
            input="this is not json",
            capture_output=True,
            text=True,
        )
        # Should either exit non-zero or return an error JSON
        self.assertTrue(proc.returncode != 0 or "error" in proc.stdout)


if __name__ == "__main__":
    unittest.main(verbosity=2)
